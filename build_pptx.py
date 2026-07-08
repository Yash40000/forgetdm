#!/usr/bin/env python3
"""Generate ForgeTDM_Demo_Deck.pptx (native, editable) — mirrors the HTML demo deck.
Run:  py -m pip install python-pptx   then   py build_pptx.py
Output: ForgeTDM_Demo_Deck.pptx next to this script."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
NAVY   = RGBColor(0x0B, 0x25, 0x45)
NAVY2  = RGBColor(0x13, 0x31, 0x5C)
BLUE   = RGBColor(0x1D, 0x6F, 0xB8)
BLUE2  = RGBColor(0x2E, 0x86, 0xDE)
ICE    = RGBColor(0xE8, 0xF0, 0xFE)
ICE2   = RGBColor(0xF4, 0xF8, 0xFF)
INK    = RGBColor(0x15, 0x21, 0x2B)
MUTED  = RGBColor(0x5C, 0x6B, 0x7B)
LINE   = RGBColor(0xDC, 0xE4, 0xED)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ICEBLU = RGBColor(0x9F, 0xC2, 0xEE)
SUBBL  = RGBColor(0xCB, 0xDD, 0xF6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, radius=False, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if shadow:
        el = shp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {}); el.append(ef)
        sh = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad': '90000', 'dist': '30000', 'dir': '5400000', 'rotWithShape': '0'})
        ef.append(sh)
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '0B2545'}); sh.append(clr)
        al = clr.makeelement(qn('a:alpha'), {'val': '14000'}); clr.append(al)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for t, size, color, bold, *rest in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
            r.font.name = rest[0] if rest else 'Calibri'
    return tb


def header(s, eyebrow, title, num):
    text(s, 0.73, 0.62, 9, 0.3, [[(eyebrow, 12, BLUE, True)]])
    text(s, 0.73, 0.9, 10, 0.7, [[(title, 30, NAVY, True, 'Cambria')]])
    text(s, 12.0, 0.95, 0.7, 0.4, [[(num, 14, BLUE, True, 'Cambria')]], align=PP_ALIGN.RIGHT)
    box(s, 0.73, 1.62, 11.85, 0.028, fill=LINE)


def card(s, x, y, w, h, title, body, tint=False, tag=None):
    box(s, x, y, w, h, fill=(ICE2 if tint else WHITE), line=LINE, radius=True, shadow=not tint)
    ty = y + 0.18
    if tag:
        text(s, x + 0.22, ty, w - 0.4, 0.25, [[(tag, 9.5, BLUE, True)]]); ty += 0.3
    text(s, x + 0.22, ty, w - 0.44, 0.4, [[(title, 14.5, NAVY, True)]])
    text(s, x + 0.22, ty + 0.42, w - 0.44, h - (ty - y) - 0.5, [[(body, 11.5, MUTED, False)]])


# ---------- 1: COVER ----------
s = slide(NAVY)
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
box(s, 0.73, 1.85, 0.62, 0.62, fill=BLUE2, radius=True)
text(s, 0.73, 1.9, 0.62, 0.55, [[("F", 30, WHITE, True, 'Cambria')]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 1.5, 1.86, 8, 0.7, [[("ForgeTDM", 26, WHITE, True, 'Cambria')],
                            [("TEST DATA PLATFORM", 11, ICEBLU, True)]], space=0)
text(s, 0.73, 2.95, 6, 0.3, [[("PRODUCT DEMO", 13, BLUE2, True)]])
text(s, 0.73, 3.3, 11.4, 1.7,
     [[("Enterprise Test Data Management — discover, mask, subset, synthesize, deliver", 38, WHITE, True, 'Cambria')]])
text(s, 0.73, 5.1, 10.5, 0.8,
     [[("One governed platform for safe, realistic, on-demand test data across relational and mainframe systems.", 19, SUBBL, False)]])
cx = 0.73
for label in ["Compliant masking", "Referential synthetic data", "Mainframe-native", "AI Copilot"]:
    w = 0.35 + len(label) * 0.115
    b = box(s, cx, 6.05, w, 0.42, fill=NAVY2, line=ICEBLU, radius=True)
    text(s, cx, 6.05, w, 0.42, [[(label, 13, RGBColor(0xDC, 0xE9, 0xFA), True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    cx += w + 0.18

# ---------- 2: CAPABILITIES ----------
s = slide()
header(s, "WHAT IT DOES", "Capabilities", "02")
caps = [
    ("PII Discovery", "Dual-signal scan of names + sampled values; classify, approve, auto-build a policy."),
    ("Masking", "Deterministic, format-preserving, conditional (WHERE), and in-place at billion-row scale."),
    ("DataScope & Subsetting", "Blueprint-driven, right-sized subsets with full referential integrity + overrides."),
    ("Synthetic Data", "Multi-table, FK-aware with ER view, type-safe generators; load to DB or CSV/JSON/SQL."),
    ("Mainframe", "COBOL copybooks, EBCDIC & COMP-3 — byte-faithful masking, Zowe delivery."),
    ("Virtualization", "Space-efficient snapshot & clone — stand up masked environments fast."),
    ("Validation", "Automated leak, format, referential and domain checks prove data is safe and valid."),
    ("AI Copilot", "Pluggable LLM assistant that reads your catalog and runs work — with approval gates."),
]
gx, gy, gw, gh, gp = 0.73, 1.95, 2.86, 2.15, 0.16
for i, (t, b) in enumerate(caps):
    col, row = i % 4, i // 4
    card(s, gx + col * (gw + gp), gy + row * (gh + 0.18), gw, gh, t, b)

# ---------- 3: USE CASES ----------
s = slide()
header(s, "WHERE IT'S USED", "Use cases", "03")
uses = [
    ("COMPLIANCE", "Safe non-production", "Remove PII from dev/test/QA with consistent, evidenced masking — GDPR, HIPAA, PCI."),
    ("RIGHT-SIZE", "Referential subsets", "Carve a small, intact slice of production by a driver + filter, FK closure preserved."),
    ("GREENFIELD", "Synthetic environments", "Generate realistic, related data for a brand-new environment — no source needed."),
    ("SCALE", "In-place masking", "Mask huge production-clone tables where source = target, chunked and restartable."),
    ("LEGACY", "Mainframe masking", "Mask fixed/variable EBCDIC datasets from copybooks without leaving the platform."),
    ("QA VELOCITY", "On-demand & defect repro", "Self-service data and compliant reproductions of production issues, fast."),
]
gx, gy, gw, gh, gp = 0.73, 1.95, 3.86, 2.15, 0.18
for i, (tag, t, b) in enumerate(uses):
    col, row = i % 3, i // 3
    card(s, gx + col * (gw + gp), gy + row * (gh + 0.22), gw, gh, t, b, tint=True, tag=tag)

# ---------- 4: TECHNICAL SPEC ----------
s = slide()
header(s, "UNDER THE HOOD", "Technical specification", "04")
spec = [
    ("Platform", "Java 17 · Spring Boot 3 service, REST API, browser UI; Flyway-managed schema."),
    ("Data engines", "Postgres, Oracle, SQL Server, DB2, MySQL, H2 — plus mainframe files via copybook codecs."),
    ("Masking engine", "Deterministic HMAC — same input to same output, irreversible; format-preserving; conditional."),
    ("Scale", "Keyset-paginated, per-chunk committed in-place updates via staging join; parallel workers."),
    ("Deployment", "Docker / containers, on-prem or in-VPC; metadata in Postgres; env-driven config."),
    ("AI integration", "Any OpenAI-compatible endpoint (cloud or local) — metadata only, never raw PII."),
    ("Governance", "Full audit trail, human approval gates, reproducible runs, role-ready access."),
    ("Integrity", "Auto FK detection, parent-first load order, type-safe generation & length-aware loads."),
]
gy = 1.95
for i, (k, v) in enumerate(spec):
    col, row = i % 2, i // 2
    x = 0.73 + col * 6.0
    y = gy + row * 0.92
    text(s, x, y, 1.6, 0.4, [[(k, 14.5, NAVY, True)]])
    text(s, x + 1.65, y, 4.15, 0.85, [[(v, 13, INK, False)]])
by = gy + 4 * 0.92 + 0.15
bx = 0.73
for label in ["GDPR", "HIPAA", "PCI-DSS", "REST API / CI-CD ready", "On-prem / VPC"]:
    w = 0.4 + len(label) * 0.105
    box(s, bx, by, w, 0.44, fill=WHITE, line=LINE, radius=True)
    text(s, bx, by, w, 0.44, [[(label, 13.5, NAVY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx += w + 0.18

# ---------- 5: BENEFITS ----------
s = slide()
header(s, "WHY IT MATTERS", "Benefits", "05")
stats = [("Days -> Hrs", "Provisioning time, via self-service & chunked delivery."),
         ("1 platform", "Replaces scattered scripts across relational & mainframe."),
         ("100% audit", "Every masking & AI action logged and explainable."),
         ("Lower risk", "Consistent, evidenced removal of PII from non-prod.")]
gx, gw, gp = 0.73, 2.86, 0.16
for i, (big, lab) in enumerate(stats):
    x = gx + i * (gw + gp)
    box(s, x, 1.95, gw, 1.5, fill=ICE2, line=LINE, radius=True)
    text(s, x + 0.22, 2.1, gw - 0.44, 0.6, [[(big, 26, NAVY, True, 'Cambria')]])
    text(s, x + 0.22, 2.78, gw - 0.44, 0.6, [[(lab, 12, MUTED, False)]])
bens = [("Compliance by design", "Deterministic masking and an audit trail give auditors the evidence pack on demand."),
        ("Faster, self-service delivery", "Teams get right-sized, masked or synthetic data without DBA tickets."),
        ("Trustworthy data", "Referential integrity, type-safety and validation mean data that works in test.")]
gx, gw, gp = 0.73, 3.86, 0.18
for i, (t, b) in enumerate(bens):
    card(s, gx + i * (gw + gp), 3.75, gw, 1.7, t, b)
text(s, 0.73, 5.65, 11, 0.3, [[("Figures are representative value directions, quantified per environment during a pilot.", 11, MUTED, False, 'Calibri')]])

# ---------- 6: FUTURE ENHANCEMENTS ----------
s = slide()
header(s, "WHAT'S NEXT", "Future enhancements", "06")
phases = [
    ("NEAR TERM", ["Shadow-table rebuild for in-place masking of unique-indexed columns",
                   "Richer generators (currency/code, address packs) with smart suggestions",
                   "Load-time 'uncovered NOT NULL column' warnings"]),
    ("MID TERM", ["Agentic provisioning: conversational, end-to-end with approvals",
                  "Self-healing validation that explains & proposes fixes",
                  "More connectors (NoSQL, message queues, cloud warehouses)"]),
    ("STRATEGIC", ["Entity-based delivery & refresh, K2View-style",
                   "Test-case-driven coverage generation",
                   "Scheduling, approvals & environment lifecycle automation"]),
]
gx, gw, gp = 0.73, 3.86, 0.18
for i, (ph, items) in enumerate(phases):
    x = gx + i * (gw + gp)
    box(s, x, 1.95, gw, 3.5, fill=ICE2, line=LINE, radius=True)
    text(s, x + 0.24, 2.18, gw - 0.5, 0.3, [[(ph, 12, BLUE, True)]])
    paras = [[("•  " + it, 13.5, INK, False)] for it in items]
    text(s, x + 0.24, 2.62, gw - 0.5, 2.6, paras, space=8)

# ---------- 7: CLOSING ----------
s = slide(NAVY)
box(s, 0, 0, 13.333, 7.5, fill=NAVY)
text(s, 0.73, 1.55, 8, 0.3, [[("LET'S SEE IT LIVE", 13, RGBColor(0x7F, 0xB2, 0xF0), True)]])
text(s, 0.73, 1.95, 11.4, 1.5,
     [[("Safe, realistic, on-demand test data — from one governed platform", 36, WHITE, True, 'Cambria')]])
text(s, 0.73, 3.35, 10.5, 0.7,
     [[("Relational and mainframe, masked and synthetic, AI-accelerated and compliance-ready.", 18, SUBBL, False)]])
steps = [("1", "Discover & mask sensitive data with deterministic, evidenced rules."),
         ("2", "Subset or synthesize referentially-intact datasets on demand."),
         ("3", "Deliver to any target — relational, files, or mainframe."),
         ("4", "Govern every run with audit, approvals and validation.")]
for i, (n, t) in enumerate(steps):
    col, row = i % 2, i // 2
    x = 0.73 + col * 5.7
    y = 4.35 + row * 0.7
    box(s, x, y, 0.42, 0.42, fill=NAVY2, line=ICEBLU, radius=True)
    text(s, x, y, 0.42, 0.42, [[(n, 14, ICEBLU, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.6, y - 0.02, 4.9, 0.6, [[(t, 14.5, RGBColor(0xD9, 0xE6, 0xF8), False)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.73, 6.35, 11, 0.4, [[("Next step — a scoped pilot on one application.", 16, ICEBLU, True)]])

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ForgeTDM_Demo_Deck.pptx")
prs.save(out)
print("Saved:", out)
